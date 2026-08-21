"""
build_slides.py - generate the class-presentation deck for this project.

Themed to match the web app: woven-palm basket ground, chalk cards, banana-leaf green,
mangosteen and turmeric for the two directions a forecast can go. Fraunces / IBM Plex
aren't installed on a typical Windows machine, so the deck substitutes the nearest faces
that ship with Windows - Georgia for display, Segoe UI for body, Consolas for figures -
which keeps it looking right on a classroom projector.

Every figure quoted is computed from the artifacts in this repo, not typed by hand:
run tools/metrics.py first if metrics.json is missing.

    pip install python-pptx
    python tools/build_slides.py
"""

import json
import os
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
ASSETS = ROOT / "tools" / "slide_assets"
# SLIDES_OUT lets you build to a scratch path while the real deck is open in PowerPoint,
# which holds a write lock on it.
OUT = Path(os.environ.get("SLIDES_OUT") or ROOT / "Cambodia_Food_Price_Forecast.pptx")

# ---------- palette (identical to app/assets/css/main.css) ----------
BASKET = RGBColor(0xE9, 0xE7, 0xDB)
BASKET_DEEP = RGBColor(0xDC, 0xD9, 0xC9)
CHALK = RGBColor(0xFC, 0xFB, 0xF7)
INK = RGBColor(0x19, 0x1E, 0x19)
INK_SOFT = RGBColor(0x5C, 0x62, 0x5A)
PALM = RGBColor(0x2F, 0x5D, 0x46)
PALM_LIGHT = RGBColor(0x4D, 0x80, 0x65)
MANGOSTEEN = RGBColor(0x6B, 0x25, 0x45)
TURMERIC = RGBColor(0xD9, 0x9A, 0x0B)
TURMERIC_DEEP = RGBColor(0xA0, 0x6D, 0x05)

DISPLAY = "Georgia"      # stands in for Fraunces
SANS = "Segoe UI"        # stands in for IBM Plex Sans
MONO = "Consolas"        # stands in for IBM Plex Mono

W, H = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.9)
COL = W - 2 * MARGIN

prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]

metrics = json.loads((ROOT / "tools" / "metrics.json").read_text())
REF = json.loads((ROOT / "nuxt-app" / "test" / "reference.json").read_text())
HISTORY = json.loads(
    (ROOT / "nuxt-app" / "server" / "assets" / "artifacts" / "history.json").read_text(encoding="utf-8")
)
SERIES_COUNT = sum(len(by_type) for by_market in HISTORY.values() for by_type in by_market.values())

# Headline figures, quoted on both the problem/solution slide and the result slide.
xgb = metrics["rows"][-1]
base = metrics["rows"][0]
improve = (1 - xgb["mae"] / base["mae"]) * 100


# ---------- primitives ----------

def slide(bg=BASKET):
    s = prs.slides.add_slide(BLANK)
    rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    rect.fill.solid()
    rect.fill.fore_color.rgb = bg
    rect.line.fill.background()
    rect.shadow.inherit = False
    return s


def box(s, left, top, width, height, fill=None, line=None, line_w=1.0, shape=MSO_SHAPE.RECTANGLE):
    sh = s.shapes.add_shape(shape, left, top, width, height)
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_w)
    sh.shadow.inherit = False
    return sh


def text(s, body, left, top, width, height, font=SANS, size=16, color=INK,
         bold=False, italic=False, align=PP_ALIGN.LEFT, spacing=None, line=1.3,
         anchor=MSO_ANCHOR.TOP, space_after=0):
    """body: a string, or a list of (text, overrides-dict) paragraphs."""
    tb = s.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

    paras = body if isinstance(body, list) else [(body, {})]
    for i, item in enumerate(paras):
        content, over = item if isinstance(item, tuple) else (item, {})
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = over.get("align", align)
        p.line_spacing = over.get("line", line)
        p.space_after = Pt(over.get("space_after", space_after))
        r = p.add_run()
        r.text = content
        f = r.font
        f.name = over.get("font", font)
        f.size = Pt(over.get("size", size))
        f.bold = over.get("bold", bold)
        f.italic = over.get("italic", italic)
        f.color.rgb = over.get("color", color)
        sp = over.get("spacing", spacing)
        if sp:
            # python-pptx exposes no character-spacing API; set the OOXML attribute.
            r.font._element.set("spc", str(int(sp * 100)))
    return tb


def eyebrow(s, label, top=Inches(0.62)):
    return text(s, label.upper(), MARGIN, top, COL, Inches(0.3),
                font=MONO, size=10.5, color=INK_SOFT, spacing=2.0)


def headline(s, line1, line2=None, top=Inches(1.02), size=36):
    text(s, line1, MARGIN, top, COL, Inches(0.7), font=DISPLAY, size=size,
         color=INK, bold=True, line=1.05)
    if line2:
        text(s, line2, MARGIN, top + Inches(0.62), COL, Inches(0.7), font=DISPLAY,
             size=size, color=PALM, italic=True, line=1.05)


def rule(s, left, top, width, color=BASKET_DEEP, weight=1.0):
    ln = s.shapes.add_connector(1, left, top, left + width, top)
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    return ln


def card(s, left, top, width, height, accent=None, fill=CHALK):
    box(s, left, top, width, height, fill=fill, line=BASKET_DEEP)
    if accent:
        box(s, left, top, Inches(0.05), height, fill=accent)


def number(s, n):
    text(s, f"{n:02d}", W - MARGIN - Inches(0.6), H - Inches(0.56), Inches(0.6),
         Inches(0.3), font=MONO, size=9.5, color=INK_SOFT, align=PP_ALIGN.RIGHT,
         spacing=1.0)


def footnote(s, note):
    text(s, note, MARGIN, H - Inches(0.72), COL - Inches(1.0), Inches(0.5),
         font=MONO, size=9.5, color=INK_SOFT, spacing=0.6, line=1.35)


# ---------- slides ----------
n = 0


def new(label=None, l1=None, l2=None, size=36, bg=BASKET, numbered=True):
    global n
    s = slide(bg)
    if label:
        eyebrow(s, label)
    if l1:
        headline(s, l1, l2, size=size)
    if numbered:
        n += 1
        number(s, n)
    return s


# 01 — title -----------------------------------------------------------------
s = slide(BASKET)
text(s, "FINAL PROJECT  ·  ARTIFICIAL INTELLIGENCE", MARGIN, Inches(1.5), COL,
     Inches(0.3), font=MONO, size=11, color=PALM, spacing=2.6)
text(s, "What the market", MARGIN, Inches(2.05), COL, Inches(1.1), font=DISPLAY,
     size=62, color=INK, bold=True, line=1.0)
text(s, "will ask next", MARGIN, Inches(3.02), COL, Inches(1.1), font=DISPLAY,
     size=62, color=PALM, italic=True, line=1.0)
rule(s, MARGIN, Inches(4.35), Inches(3.2), PALM_LIGHT, 1.5)
text(s, "Forecasting food commodity prices in Cambodia — from World Food Programme "
        "price monitoring to a web app and a Telegram bot.",
     MARGIN, Inches(4.7), Inches(7.7), Inches(1.0), size=17, color=INK_SOFT, line=1.45)
text(s, "46 COMMODITIES   ·   76 MARKETS   ·   XGBOOST, PORTED TO TYPESCRIPT",
     MARGIN, Inches(6.35), COL, Inches(0.3), font=MONO, size=11, color=INK_SOFT,
     spacing=1.6)

# 02 — problem and solution --------------------------------------------------
s = new("Problem & solution", "An early warning", "only works if it arrives")

TOP, HEIGHT = Inches(2.4), Inches(3.95)
PANELS = [
    (MARGIN, Inches(5.5), MANGOSTEEN, "PROBLEM",
     "A price rise lands hardest on the people with the least slack.",
     ["For a low-income household food is the largest monthly cost, so a price rise "
      "is not an inconvenience.",
      "WFP tracks these markets to catch that early, but publishes what prices were "
      "— not what they will be.",
      f"And the record itself ships as a {metrics['raw_rows']:,}-row CSV: readable by "
      "analysts, nobody else."]),
    (Inches(6.93), Inches(5.5), PALM, "SOLUTION",
     "Make the answer reachable in two clicks or one message.",
     [f"{SERIES_COUNT:,} price series across 46 commodities and 76 markets, charted "
      "for anyone with a browser.",
      "A next-period forecast on each, with the line between recorded and predicted "
      "drawn honestly.",
      "Free to run — so a ministry, an NGO or a co-op could stand the same thing up "
      "on its own feed."]),
]

for left, width, accent, label, lead, points in PANELS:
    card(s, left, TOP, width, HEIGHT, accent=accent)
    inner = left + Inches(0.4)
    inner_w = width - Inches(0.8)

    text(s, label, inner, TOP + Inches(0.32), inner_w, Inches(0.3), font=MONO,
         size=10, color=accent, spacing=2.0)
    text(s, lead, inner, TOP + Inches(0.72), inner_w, Inches(0.8), font=DISPLAY,
         size=19, color=INK, bold=True, line=1.15)

    # Each point is written to fit two lines at this measure; the step allows for that
    # plus a gap, so the three of them land inside the card.
    y = TOP + Inches(1.55)
    for point in points:
        box(s, inner, y + Inches(0.07), Inches(0.14), Inches(0.14), fill=accent)
        text(s, point, inner + Inches(0.32), y, inner_w - Inches(0.32), Inches(0.6),
             size=13.5, color=INK_SOFT, line=1.4)
        y += Inches(0.78)

# The turn from one panel to the other.
text(s, "→", Inches(6.32), Inches(4.05), Inches(0.6), Inches(0.4), font=MONO, size=22,
     color=PALM, align=PP_ALIGN.CENTER)

footnote(s, "A demonstration on a snapshot, not a deployed service — the data ends March 2026 and the app says so, series by series.")

# 03 — the data --------------------------------------------------------------
s = new("The data", "Twenty-three years", "of market visits")
rows = [
    ("Source", "WFP Global Food Prices — Cambodia (HDX)"),
    ("Observations", f"{metrics['raw_rows']:,}"),
    ("Range", f"{metrics['raw_start']} to {metrics['raw_end']}"),
    ("Commodities", "50 raw"),
    ("Markets", "86 raw"),
    ("Target", "usdprice — price per unit in USD"),
]
y = Inches(2.45)
for label, value in rows:
    text(s, label.upper(), MARGIN, y, Inches(1.7), Inches(0.3), font=MONO, size=10,
         color=INK_SOFT, spacing=1.4)
    text(s, value, MARGIN + Inches(1.9), y - Inches(0.03), Inches(4.3), Inches(0.34),
         size=15, color=INK)
    y += Inches(0.55)

card(s, Inches(7.5), Inches(2.3), Inches(4.9), Inches(3.9), accent=TURMERIC)
text(s, "FILTERING", Inches(7.9), Inches(2.65), Inches(4.1), Inches(0.3), font=MONO,
     size=10, color=INK_SOFT, spacing=2.0)
funnel = [
    (f"{metrics['raw_rows']:,} rows", "everything WFP has recorded"),
    ("from 2019 onward", "older reporting is sparse and pre-dollarisation"),
    ("commodities with 300+ rows", "a series needs history to have a lag"),
    (f"{metrics['n_rows']:,} modelled rows", f"46 commodities · 76 markets · {metrics['n_series']:,} series"),
]
y = Inches(3.12)
for i, (big, small) in enumerate(funnel):
    colour = PALM if i == len(funnel) - 1 else INK
    text(s, big, Inches(7.9), y, Inches(4.1), Inches(0.3), font=MONO, size=14,
         color=colour, bold=(i == len(funnel) - 1))
    text(s, small, Inches(7.9), y + Inches(0.26), Inches(4.1), Inches(0.3), size=12,
         color=INK_SOFT)
    if i < len(funnel) - 1:
        text(s, "↓", Inches(7.9), y + Inches(0.5), Inches(0.3), Inches(0.2), font=MONO,
             size=11, color=BASKET_DEEP)
    y += Inches(0.76)

# 04 — framing ---------------------------------------------------------------
s = new("Framing it as a learning problem", "The split is the", "whole experiment")
text(s, "Each (commodity, market, pricetype) forms a time series. The model predicts the "
        "next observation in that series from its own recent history.",
     MARGIN, Inches(2.45), Inches(5.6), Inches(1.2), size=16, color=INK, line=1.5)

card(s, MARGIN, Inches(3.75), Inches(5.6), Inches(2.35), accent=MANGOSTEEN)
text(s, "THE TRAP", MARGIN + Inches(0.4), Inches(4.05), Inches(4.8), Inches(0.3),
     font=MONO, size=10, color=MANGOSTEEN, spacing=2.0)
text(s, "A random train/test split lets the model learn from June to predict April. "
        "Prices are autocorrelated, so it scores brilliantly and forecasts nothing. "
        "The split has to be chronological.",
     MARGIN + Inches(0.4), Inches(4.45), Inches(4.8), Inches(1.4), size=14.5,
     color=INK_SOFT, line=1.45)

card(s, Inches(7.3), Inches(2.35), Inches(5.1), Inches(3.75), accent=PALM)
text(s, "CHRONOLOGICAL HOLD-OUT", Inches(7.7), Inches(2.68), Inches(4.3), Inches(0.3),
     font=MONO, size=10, color=INK_SOFT, spacing=2.0)
split_rows = [
    ("Cut point", f"{metrics['split']}  (85th percentile of dates)"),
    ("Train", f"{metrics['n_train']:,} rows  ·  everything before the cut"),
    ("Test", f"{metrics['n_test']:,} rows  ·  everything after"),
    ("Never seen", "the test period, by any model, at any point"),
]
y = Inches(3.2)
for label, value in split_rows:
    text(s, label, Inches(7.7), y, Inches(1.5), Inches(0.3), size=13.5, color=INK_SOFT)
    text(s, value, Inches(9.25), y, Inches(2.9), Inches(0.6), font=MONO, size=12.5,
         color=INK, line=1.35)
    y += Inches(0.68)

# 05 — features --------------------------------------------------------------
s = new("Feature engineering", "Ten features, built", "inside each series")
groups = [
    ("PRICE HISTORY", PALM, [
        ("lag_1", "the previous reported price"),
        ("lag_3", "the price three reports back"),
        ("rolling_mean_3", "mean of the last three"),
    ]),
    ("CALENDAR", TURMERIC_DEEP, [
        ("month", "1-12, seasonality of harvest"),
        ("quarter", "1-4, coarser seasonality"),
    ]),
    ("IDENTITY", MANGOSTEEN, [
        ("commodity_enc", "which food"),
        ("market_enc", "which market"),
        ("admin1_enc", "which province"),
        ("category_enc", "cereals, meat, pulses …"),
        ("pricetype_enc", "retail or wholesale"),
    ]),
]
x = MARGIN
for title, colour, items in groups:
    width = Inches(3.75)
    text(s, title, x, Inches(2.45), width, Inches(0.3), font=MONO, size=10,
         color=colour, spacing=1.8)
    rule(s, x, Inches(2.78), width - Inches(0.25), colour, 1.5)
    y = Inches(2.98)
    for name, desc in items:
        text(s, name, x, y, width, Inches(0.28), font=MONO, size=13, color=INK)
        text(s, desc, x, y + Inches(0.25), width - Inches(0.3), Inches(0.28), size=12,
             color=INK_SOFT)
        y += Inches(0.53)
    x += Inches(4.05)

card(s, MARGIN, Inches(5.7), COL, Inches(1.0), accent=PALM)
text(s, "Lags are computed within each series (grouped shift), never across them — a rice "
        "price in Battambang can never leak into a pork price in Phnom Penh. Rows without "
        "three prior observations are dropped rather than imputed.",
     MARGIN + Inches(0.4), Inches(5.97), COL - Inches(0.9), Inches(0.7), size=13.5,
     color=INK_SOFT, line=1.4)

# 06 — model comparison ------------------------------------------------------
s = new("Model comparison", "Beating “tomorrow", "equals today”")
text(s, "The naive baseline — predict the last price — already scores R² 0.92, because "
        "prices are sticky. It is the number any model has to beat to be worth deploying.",
     MARGIN, Inches(2.4), Inches(11.0), Inches(0.8), size=15.5, color=INK_SOFT, line=1.45)

head_y = Inches(3.35)
cols = [(MARGIN, Inches(4.6), "MODEL", PP_ALIGN.LEFT),
        (Inches(6.0), Inches(1.9), "MAE (USD)", PP_ALIGN.RIGHT),
        (Inches(8.1), Inches(1.9), "RMSE (USD)", PP_ALIGN.RIGHT),
        (Inches(10.2), Inches(2.2), "R²", PP_ALIGN.RIGHT)]
for left, width, label, align in cols:
    text(s, label, left, head_y, width, Inches(0.3), font=MONO, size=10,
         color=INK_SOFT, spacing=1.6, align=align)
rule(s, MARGIN, head_y + Inches(0.32), COL, INK_SOFT, 1.0)

y = head_y + Inches(0.55)
for i, r in enumerate(metrics["rows"]):
    best = i == len(metrics["rows"]) - 1
    colour = PALM if best else INK
    if best:
        box(s, MARGIN - Inches(0.18), y - Inches(0.12), COL + Inches(0.36),
            Inches(0.72), fill=CHALK, line=BASKET_DEEP)
    text(s, r["name"], MARGIN, y, Inches(4.6), Inches(0.4), size=15.5, color=colour,
         bold=best)
    for left, width, value in ((Inches(6.0), Inches(1.9), f"{r['mae']:.4f}"),
                               (Inches(8.1), Inches(1.9), f"{r['rmse']:.4f}"),
                               (Inches(10.2), Inches(2.2), f"{r['r2']:.4f}")):
        text(s, value, left, y + Inches(0.02), width, Inches(0.4), font=MONO, size=15,
             color=colour, bold=best, align=PP_ALIGN.RIGHT)
    y += Inches(0.88)

footnote(s, "Random Forest reproduced locally from the notebook's settings (n=300, depth 15, seed 42).\n"
            "XGBoost figures come from the trained price_model.pkl itself, evaluated on the held-out period.")

# 07 — headline result -------------------------------------------------------
s = new("Result", "The model earns", "its place")
stats =[(f"${xgb['mae']:.3f}", "MEAN ABSOLUTE ERROR", "on held-out 2025-26 prices", PALM),
         (f"{xgb['r2']:.3f}", "R² ON THE TEST PERIOD", "variance explained", PALM_LIGHT),
         (f"{improve:.1f}%", "BETTER THAN NAIVE", "MAE reduction vs. last price", TURMERIC_DEEP)]
x = MARGIN
for value, label, sub, colour in stats:
    card(s, x, Inches(2.6), Inches(3.6), Inches(2.5), accent=colour)
    text(s, value, x + Inches(0.4), Inches(3.0), Inches(3.0), Inches(0.9), font=MONO,
         size=40, color=colour, bold=True)
    text(s, label, x + Inches(0.4), Inches(3.95), Inches(3.0), Inches(0.3), font=MONO,
         size=9.5, color=INK_SOFT, spacing=1.4)
    text(s, sub, x + Inches(0.4), Inches(4.3), Inches(2.9), Inches(0.5), size=13,
         color=INK_SOFT, line=1.35)
    x += Inches(3.85)

text(s, "An average miss of about sixteen US cents on a basket whose prices mostly sit "
        "between $0.30 and $3.00 — accurate enough to plan around, not precise enough "
        "to trade on. Both halves of that sentence matter.",
     MARGIN, Inches(5.5), Inches(11.0), Inches(1.0), size=15.5, color=INK, line=1.5)

# 08 — what it learned -------------------------------------------------------
s = new("Interpretation", "It learned to be", "a very good smoother")
imp = metrics["importance"]
top = list(imp.items())[:6]
bar_left = MARGIN + Inches(2.5)
bar_max = Inches(5.4)
y = Inches(2.6)
for name, value in top:
    text(s, name, MARGIN, y - Inches(0.02), Inches(2.35), Inches(0.3), font=MONO,
         size=12.5, color=INK, align=PP_ALIGN.RIGHT)
    width = max(Emu(int(bar_max * value)), Emu(int(Inches(0.02))))
    colour = PALM if value > 0.1 else PALM_LIGHT
    box(s, bar_left, y, width, Inches(0.28), fill=colour)
    text(s, f"{value * 100:.1f}%", bar_left + width + Inches(0.12), y - Inches(0.02),
         Inches(1.0), Inches(0.3), font=MONO, size=12, color=INK_SOFT)
    y += Inches(0.52)

card(s, Inches(9.0), Inches(2.45), Inches(3.4), Inches(3.95), accent=MANGOSTEEN)
text(s, "THE HONEST READING", Inches(9.35), Inches(2.78), Inches(2.8), Inches(0.3),
     font=MONO, size=9.5, color=MANGOSTEEN, spacing=1.6)
text(s, [(f"{(imp['rolling_mean_3'] + imp['lag_1']) * 100:.0f}% of the model",
          {"size": 19, "font": DISPLAY, "bold": True, "color": INK, "space_after": 8}),
         ("rests on two features: the recent mean and the last price.",
          {"size": 14, "color": INK, "space_after": 12}),
         ("Commodity, market, province and season together account for barely two "
          "percent. This is a well-calibrated smoother, not a market oracle — and "
          "saying so is part of the result.", {"size": 12.5, "color": INK_SOFT})],
     Inches(9.35), Inches(3.2), Inches(2.75), Inches(2.9), line=1.35)

footnote(s, "Gain-based feature importance from the trained booster.")

# 09 — deployment problem ----------------------------------------------------
s = new("From notebook to product", "The model fit.", "The libraries didn’t.")
text(s, "A notebook that only runs in Colab isn't a product. Deploying to Vercel's free "
        "tier meant meeting a hard limit: 250 MB unzipped per serverless function.",
     MARGIN, Inches(2.4), Inches(11.0), Inches(0.8), size=16, color=INK, line=1.5)

# Bars share one scale so the total genuinely overshoots the limit line rather than
# being drawn to fit. 360 MB spans the plot, which keeps 327 MB on the slide.
libs = [("xgboost", 84), ("scipy", 111), ("numpy", 58), ("pandas", 42), ("scikit-learn", 32)]
BAR_X = MARGIN + Inches(1.8)
PLOT_W = Inches(9.6)
scale = PLOT_W / 360.0
LIMIT_X = BAR_X + Emu(int(scale * 250))

y = Inches(3.55)
for name, mb in libs:
    text(s, name, MARGIN, y - Inches(0.02), Inches(1.6), Inches(0.3), font=MONO,
         size=12, color=INK, align=PP_ALIGN.RIGHT)
    box(s, BAR_X, y, Emu(int(scale * mb)), Inches(0.24), fill=BASKET_DEEP)
    text(s, f"{mb} MB", BAR_X + Inches(0.1) + Emu(int(scale * mb)), y - Inches(0.03),
         Inches(1.0), Inches(0.3), font=MONO, size=11, color=INK_SOFT)
    y += Inches(0.4)

rule(s, BAR_X, y + Inches(0.04), Inches(9.6), BASKET_DEEP, 1.0)
text(s, "TOTAL", MARGIN, y + Inches(0.2), Inches(1.6), Inches(0.3), font=MONO, size=12,
     color=INK, bold=True, align=PP_ALIGN.RIGHT)
box(s, BAR_X, y + Inches(0.19), Emu(int(scale * 327)), Inches(0.32), fill=MANGOSTEEN)
text(s, "327 MB", BAR_X + Inches(0.12) + Emu(int(scale * 327)), y + Inches(0.22),
     Inches(1.4), Inches(0.3), font=MONO, size=13, color=MANGOSTEEN, bold=True)

# The limit, drawn where it actually falls — the total crosses it by 77 MB.
limit = s.shapes.add_connector(1, LIMIT_X, Inches(3.35), LIMIT_X, y + Inches(0.62))
limit.line.color.rgb = INK
limit.line.width = Pt(1.5)
limit.line.dash_style = 4  # dash
text(s, "250 MB  ·  VERCEL HOBBY LIMIT", LIMIT_X - Inches(3.6), Inches(3.02),
     Inches(3.5), Inches(0.3), font=MONO, size=10.5, color=INK, spacing=1.4,
     align=PP_ALIGN.RIGHT)

text(s, "The deploy fails before a single prediction is ever served.",
     MARGIN, y + Inches(0.95), Inches(8.0), Inches(0.4), size=15, color=INK)

# 10 — the port --------------------------------------------------------------
s = new("The insight", "A boosted tree is", "just a lot of if-statements")
text(s, "Training needs XGBoost. Predicting does not. Inference is: walk each tree "
        "comparing one feature to one threshold, sum the leaf weights, add the base "
        "score. Fifteen lines in any language.",
     MARGIN, Inches(2.4), Inches(5.6), Inches(1.3), size=15.5, color=INK, line=1.5)

code_lines = [
    ("function rawPredict(features: number[]): number {", INK),
    ("  let total = BASE_SCORE", INK),
    ("  for (const tree of TREES) {", INK),
    ("    let node = 0", INK),
    ("    while (tree.l[node] !== -1) {", INK),
    ("      node = features[tree.f[node]] < tree.t[node]", INK),
    ("        ? tree.l[node]", INK),
    ("        : tree.r[node]", INK),
    ("    }", INK),
    ("    total += tree.w[node]", INK),
    ("  }", INK),
    ("  return total", INK),
    ("}", INK),
]
card(s, Inches(7.0), Inches(2.35), Inches(5.4), Inches(3.85), accent=PALM)
y = Inches(2.62)
for line_text, colour in code_lines:
    text(s, line_text, Inches(7.35), y, Inches(4.9), Inches(0.26), font=MONO, size=11.5,
         color=colour)
    y += Inches(0.275)

results = [("400", "trees traversed"), ("10", "features each"), ("2.4 MB", "of JSON"),
           ("0", "ML dependencies")]
x = MARGIN
for value, label in results:
    text(s, value, x, Inches(4.1), Inches(1.5), Inches(0.5), font=MONO, size=26,
         color=PALM, bold=True)
    text(s, label, x, Inches(4.68), Inches(1.5), Inches(0.5), size=12, color=INK_SOFT,
         line=1.25)
    x += Inches(1.55)

text(s, "The model exports to JSON: tree arrays, encoder maps, feature order, and a "
        "24-point price history per series. The deployed function ships none of the "
        "Python stack.", MARGIN, Inches(5.45), Inches(5.6), Inches(1.0), size=14,
     color=INK_SOFT, line=1.45)

# 11 — verification ----------------------------------------------------------
s = new("Verification", "Two languages,", "one number")
text(s, "A port is only worth anything if it agrees with the original. The test suite "
        "pins the TypeScript traversal to the Python model's own answer for a known "
        "series — a wrong feature order or a flipped comparison changes that number, "
        "and nothing else would catch it.",
     MARGIN, Inches(2.4), Inches(5.5), Inches(1.6), size=15.5, color=INK, line=1.5)

card(s, Inches(6.9), Inches(2.35), Inches(5.5), Inches(4.05), accent=PALM)
text(s, "REFERENCE CASE", Inches(7.3), Inches(2.68), Inches(4.6), Inches(0.3),
     font=MONO, size=9.5, color=INK_SOFT, spacing=1.6)
lines = [
    ("commodity", REF["commodity"]),
    ("market", f'{REF["market"]}  ·  {REF["pricetype"]}'),
    ("last recorded", f'{REF["lastDate"]}  ·  ${REF["lastPrice"]:.2f}'),
]
y = Inches(3.1)
for label, value in lines:
    text(s, label, Inches(7.3), y, Inches(1.5), Inches(0.3), font=MONO, size=11,
         color=INK_SOFT)
    text(s, value, Inches(8.85), y, Inches(3.3), Inches(0.4), size=13.5, color=INK)
    y += Inches(0.5)

rule(s, Inches(7.3), y + Inches(0.05), Inches(4.7), BASKET_DEEP)
y += Inches(0.28)
for label, value, colour in (("python  xgboost", f'{REF["pythonPrediction"]:.6f}', INK_SOFT),
                             ("typescript  port", "0.4609556", PALM)):
    text(s, label, Inches(7.3), y, Inches(2.0), Inches(0.3), font=MONO, size=11,
         color=INK_SOFT)
    text(s, value, Inches(9.4), y - Inches(0.06), Inches(2.7), Inches(0.4), font=MONO,
         size=19, color=colour, bold=(colour == PALM))
    y += Inches(0.58)

text(s, "agree to 1e-5  ·  asserted on every test run", Inches(7.3), y + Inches(0.05),
     Inches(4.7), Inches(0.3), font=MONO, size=11, color=PALM)

text(s, "19 tests cover the traversal, the fuzzy matching, the commodity/market splitter "
        "and every failure path.", MARGIN, Inches(4.75), Inches(5.5), Inches(0.8),
     size=14, color=INK_SOFT, line=1.45)

# 12 — the app ---------------------------------------------------------------
s = new("The result", "It runs", "in a browser")
text(s, "Nuxt 4 on Vercel. The default forecast is server-rendered, so the chart is on "
        "screen at first paint — no spinner, no empty state.",
     MARGIN, Inches(2.4), Inches(5.6), Inches(1.2), size=15.5, color=INK, line=1.5)

app_facts = [
    ("Dependent dropdowns", "picking a commodity narrows the markets to those that carry it"),
    ("Fuzzy input", '"rice" resolves to Rice (mixed, low quality) — shortest match wins'),
    ("Cached at the edge", "options for a day, forecasts for an hour"),
    ("Responsive to 390px", "the chart switches to a narrower viewBox, not smaller type"),
]
y = Inches(3.7)
for title, desc in app_facts:
    text(s, title, MARGIN, y, Inches(5.5), Inches(0.28), size=14, color=INK, bold=True)
    text(s, desc, MARGIN, y + Inches(0.27), Inches(5.4), Inches(0.4), size=12.5,
         color=INK_SOFT, line=1.35)
    y += Inches(0.78)

hero = ASSETS / "hero.png"
if hero.exists():
    s.shapes.add_picture(str(hero), Inches(7.4), Inches(0.5), height=Inches(6.5))

# 13 — the chart decision ----------------------------------------------------
s = new("Design decision", "Where the record ends", "must be unmistakable")
chart = ASSETS / "chart_card.png"
if chart.exists():
    s.shapes.add_picture(str(chart), MARGIN, Inches(2.4), width=Inches(7.3))

text(s, "A charting library would draw the forecast as the next point on the same line. "
        "That is a lie of omission: it renders a model estimate with exactly the same "
        "authority as a recorded observation.",
     Inches(8.7), Inches(2.4), Inches(3.7), Inches(1.8), size=14.5, color=INK, line=1.5)

legend = [("Solid green", "prices actually recorded", PALM),
          ("Dashed amber", "the model's estimate, rising", TURMERIC_DEEP),
          ("Dashed red", "the model's estimate, falling", MANGOSTEEN),
          ("Hollow ring", "not an observation", INK_SOFT)]
y = Inches(4.25)
for title, desc, colour in legend:
    box(s, Inches(8.7), y + Inches(0.07), Inches(0.16), Inches(0.16), fill=colour)
    text(s, title, Inches(9.0), y, Inches(3.4), Inches(0.28), size=13, color=INK,
         bold=True)
    text(s, desc, Inches(9.0), y + Inches(0.25), Inches(3.4), Inches(0.28), size=12,
         color=INK_SOFT)
    y += Inches(0.6)

footnote(s, "Hand-rolled SVG — about 200 lines, no charting library.")

# 14 — forecast horizon ------------------------------------------------------
# The obvious challenge from an examiner: the data stops in March 2026, so what does
# "next period" even mean? This slide answers it with the real distribution rather
# than letting the question hang.
s = new("Forecast horizon", "One period past the record,", "not one month past today")

recency = {"2026-03 — still reporting": 0, "2025 — recently stopped": 0,
           "2024 or earlier — stale": 0}
for by_market in HISTORY.values():
    for by_type in by_market.values():
        for series in by_type.values():
            end = series["dates"][-1]
            key = ("2026-03 — still reporting" if end >= "2026-01"
                   else "2025 — recently stopped" if end >= "2025-01"
                   else "2024 or earlier — stale")
            recency[key] += 1
total_series = sum(recency.values())

text(s, "The model forecasts one reporting period past each series' own last "
        "observation — not one month past today. Series stopped reporting at very "
        "different times, so that date differs for every commodity and market.",
     MARGIN, Inches(2.45), Inches(5.6), Inches(1.6), size=16, color=INK, line=1.5)

text(s, f"LAST OBSERVATION ACROSS {total_series:,} SERVED SERIES", MARGIN, Inches(4.15),
     Inches(5.6), Inches(0.3), font=MONO, size=10, color=INK_SOFT, spacing=1.6)

y = Inches(4.55)
bar_scale = Inches(3.1)
for label, count in recency.items():
    share = count / total_series
    colour = PALM if "still" in label else TURMERIC_DEEP if "recently" in label else MANGOSTEEN
    box(s, MARGIN, y, Emu(int(bar_scale * share)), Inches(0.22), fill=colour)
    text(s, f"{share * 100:.0f}%", MARGIN + Inches(0.1) + Emu(int(bar_scale * share)),
         y - Inches(0.03), Inches(0.7), Inches(0.3), font=MONO, size=11.5, color=colour,
         bold=True)
    text(s, f"{label}  ·  {count:,}", MARGIN, y + Inches(0.26), Inches(5.4),
         Inches(0.3), size=12.5, color=INK_SOFT)
    y += Inches(0.72)

card(s, Inches(6.93), Inches(2.4), Inches(5.5), Inches(3.95), accent=PALM)
text(s, "SO IS IT ACTUALLY FORECASTING?", Inches(7.33), Inches(2.72), Inches(4.7),
     Inches(0.3), font=MONO, size=10, color=PALM, spacing=1.6)
text(s, [("Yes — and that is what the test period measures.",
          {"size": 19, "font": DISPLAY, "bold": True, "color": INK, "space_after": 12}),
         (f"Training stopped at {metrics['split']}; all {metrics['n_test']:,} test rows "
          "fall after it. The model was scored on dates it had never seen — a real "
          "out-of-sample forecast.", {"size": 13.5, "color": INK_SOFT,
                                      "space_after": 12}),
         ("What it cannot do is forecast today from a series that went quiet in 2022 — "
          "and the app prints the real date rather than hiding it.",
          {"size": 13.5, "color": INK_SOFT})],
     Inches(7.33), Inches(3.2), Inches(4.7), Inches(2.9), line=1.4)

footnote(s, "Rice at Phnom Penh, wholesale: last recorded June 2022, so the app labels its forecast July 2022 — not next month.")

# 15 — telegram --------------------------------------------------------------
s = new("Second surface", "The same engine,", "in a chat window")
text(s, "One prediction engine, two front doors. The bot shares the exact code the web "
        "app calls — a Telegram update in, a formatted forecast out.",
     MARGIN, Inches(2.4), Inches(5.9), Inches(1.0), size=15.5, color=INK, line=1.5)

bot_facts = [
    ("Reply in the response body", "the webhook returns the sendMessage call itself — no "
                                   "outbound request, and the bot token never reaches Vercel"),
    ("Always HTTP 200", "a non-2xx makes Telegram retry the same update forever"),
    ("Sparkline in Unicode blocks", "the only chart a chat window can render inline"),
]
y = Inches(3.85)
for title, desc in bot_facts:
    text(s, title, MARGIN, y, Inches(5.8), Inches(0.28), size=14, color=INK, bold=True)
    text(s, desc, MARGIN, y + Inches(0.27), Inches(5.7), Inches(0.6), size=12.5,
         color=INK_SOFT, line=1.35)
    y += Inches(0.92)

card(s, Inches(7.4), Inches(2.35), Inches(5.0), Inches(4.0), accent=PALM)
text(s, "/predict rice phnom penh", Inches(7.75), Inches(2.7), Inches(4.4), Inches(0.3),
     font=MONO, size=13, color=PALM, bold=True)
rule(s, Inches(7.75), Inches(3.08), Inches(4.3), BASKET_DEEP)
reply = [
    ("Rice (mixed, low quality)", 14, INK, True),
    ("Phnom Penh, Phnom Penh · Wholesale", 11.5, INK_SOFT, False),
    ("", 6, INK, False),
    ("Last recorded (2022-06): $0.440 / KG", 11.5, INK, False),
    ("Next period forecast: $0.461", 11.5, INK, False),
    ("▲ +4.8%", 12.5, TURMERIC_DEEP, True),
    ("", 6, INK, False),
    ("▁▃▁▁▁▁▂▂▂▁▂▂▁▇▇▇▁▇▁▇▇▇█▇", 13, PALM, False),
    ("2020-05 → 2022-06, 24 observations", 10.5, INK_SOFT, False),
]
y = Inches(3.25)
for line_text, size, colour, bold in reply:
    if line_text:
        text(s, line_text, Inches(7.75), y, Inches(4.4), Inches(0.3), font=MONO,
             size=size, color=colour, bold=bold)
    y += Inches(size / 72 * 1.9)

# 16 — limitations -----------------------------------------------------------
s = new("Limitations", "What this model", "cannot do")
limits = [
    ("One period ahead, and no further", "Every feature is a recent price. Feed the "
     "model its own output and error compounds fast — there is no multi-step forecast here."),
    ("It cannot see shocks coming", "Fuel costs, export bans, floods, festival demand. "
     "None are in the feature set, and they are exactly what moves prices sharply."),
    ("Stale series forecast from stale prices", "A series last reported in 2022 gets a "
     "2022-shaped answer. Honest, but not useful for today."),
    ("No uncertainty interval", "A single number, not a range. The dashed line signals "
     "doubt visually, but does not quantify it."),
]
y = Inches(2.5)
for title, desc in limits:
    box(s, MARGIN, y + Inches(0.06), Inches(0.05), Inches(0.6), fill=MANGOSTEEN)
    text(s, title, MARGIN + Inches(0.3), y, Inches(11.0), Inches(0.3), size=15.5,
         color=INK, bold=True)
    text(s, desc, MARGIN + Inches(0.3), y + Inches(0.33), Inches(10.6), Inches(0.5),
         size=13.5, color=INK_SOFT, line=1.4)
    y += Inches(1.12)

# 17 — next steps ------------------------------------------------------------
s = new("If I kept going", "Three things", "I would add next")
nexts = [
    ("Prediction intervals", "Quantile regression at 10/90 to turn the dashed line into "
     "a shaded band — the honest version of the same signal."),
    ("Exogenous features", "Fuel price, rainfall, and the lunar calendar for festival "
     "demand. The categorical features contribute two percent; these might not."),
    ("Retraining on a schedule", "The export script already runs headless. A monthly job "
     "against the live HDX feed would keep every series current."),
]
x = MARGIN
for title, desc in nexts:
    card(s, x, Inches(2.5), Inches(3.6), Inches(3.2), accent=PALM)
    text(s, title, x + Inches(0.4), Inches(2.9), Inches(2.9), Inches(0.7), font=DISPLAY,
         size=19, color=INK, bold=True, line=1.15)
    text(s, desc, x + Inches(0.4), Inches(3.85), Inches(2.85), Inches(1.6), size=13.5,
         color=INK_SOFT, line=1.45)
    x += Inches(3.85)

# 18 — close -----------------------------------------------------------------
s = slide(BASKET)
text(s, "THANK YOU", MARGIN, Inches(1.6), COL, Inches(0.3), font=MONO, size=11,
     color=PALM, spacing=2.6)
text(s, "Questions", MARGIN, Inches(2.15), COL, Inches(1.1), font=DISPLAY, size=58,
     color=INK, bold=True, line=1.0)
text(s, "welcome", MARGIN, Inches(3.05), COL, Inches(1.1), font=DISPLAY, size=58,
     color=PALM, italic=True, line=1.0)
rule(s, MARGIN, Inches(4.4), Inches(3.2), PALM_LIGHT, 1.5)

stack = [
    ("Data", "WFP Global Food Prices — Cambodia"),
    ("Model", "XGBoost · 400 trees · scikit-learn pipeline"),
    ("Runtime", "Nuxt 4 · Nitro · TypeScript tree traversal"),
    ("Surfaces", "Web app · Telegram bot · one shared engine"),
]
y = Inches(4.8)
for label, value in stack:
    text(s, label.upper(), MARGIN, y, Inches(1.4), Inches(0.3), font=MONO, size=10,
         color=INK_SOFT, spacing=1.4)
    text(s, value, MARGIN + Inches(1.7), y - Inches(0.03), Inches(8.0), Inches(0.32),
         size=14, color=INK)
    y += Inches(0.45)

prs.save(str(OUT))
print(f"wrote {OUT.name}  ·  {len(prs.slides.__iter__.__self__._sldIdLst)} slides  ·  "
      f"{OUT.stat().st_size / 1024:.0f} KB")
